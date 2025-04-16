from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
from base64 import b64encode
from pptx.enum.shapes import PP_PLACEHOLDER
from collections import defaultdict
from pptx.dml.color import RGBColor

import requests
import json


def load_config(config_file="config.json"):
    with open(config_file, "r") as f:
        return json.load(f)

# Get configuration
config = load_config()

# Access configuration variables
project = config["project"]
team = config["team"]
pat = config["pat"]
template_path = config["template_path"]

headers = {
    "Authorization": "Basic " + b64encode(f":{pat}".encode()).decode(),
    "Content-Type": "application/json"
}

state_colors = {
    "Active": RGBColor(0, 112, 192),     # Blue
    "New": RGBColor(255, 165, 0),        # Orange
    "Resolved": RGBColor(0, 176, 80),    # Green
    "Closed": RGBColor(0, 176, 80),      # Green
}

def get_work_item_type(work_item_id):
    # Define the URL to fetch work item details (replace with your Azure DevOps base URL)
    work_item_url = f"https://tfs-product.cmf.criticalmanufacturing.com/Products/_apis/wit/workItems/{work_item_id}?api-version=6.0"

    # Make the API call
    response = requests.get(work_item_url, headers=headers, verify=False)
    print(response)


    if response.status_code == 200:
        data = response.json()
        # Return the work item type
        return data['fields'].get('System.WorkItemType', '')
    
    return None

def get_work_item_details_w_features(ids):
    ids_str = ",".join(map(str, ids))
    details_url = f"https://tfs-product.cmf.criticalmanufacturing.com/Products/_apis/wit/workitems?ids={ids_str}&$expand=relations&api-version=7.0"

    response = requests.get(details_url, headers=headers, verify=False)

    if response.status_code == 200:
        data = response.json()
        user_stories = []
        feature_ids = set()

        for item in data['value']:
            story_id = item['id']
            title = item['fields'].get('System.Title', 'No Title')
            state = item['fields'].get('System.State', 'Unknown')
            parent_id = None

            # Look for parent (Feature) in relations
            for rel in item.get('relations', []):
                if rel['rel'] == 'System.LinkTypes.Hierarchy-Reverse':
                    parent_id = int(rel['url'].split('/')[-1])
                    feature_ids.add(parent_id)

            user_stories.append({
                'id': story_id,
                'title': title,
                'parent_id': parent_id,
                'state': state
            })

        # Fetch feature titles
        features = {}
        if feature_ids:
            feature_ids_str = ",".join(map(str, feature_ids))
            feature_url = f"https://tfs-product.cmf.criticalmanufacturing.com/Products/_apis/wit/workitems?ids={feature_ids_str}&api-version=7.0"
            feature_response = requests.get(feature_url, headers=headers, verify=False)
            if feature_response.status_code == 200:
                feature_data = feature_response.json()
                for feat in feature_data['value']:
                    features[feat['id']] = {
                        'id': feat['id'],
                        'title': feat['fields'].get('System.Title', 'Untitled Feature'),
                        'state': feat['fields'].get('System.State', 'unknown')
                    }

        # Attach feature titles to user stories
        for story in user_stories:
            story['feature'] = features.get(story['parent_id'], {'id': None, 'title': 'No Feature', 'state': 'Unknown'})

        return user_stories
    else:
        print("Failed to fetch work item details", response.status_code)
        return []

def create_feature_row(feature_id, feature_title, feature_state, text_frame):
    feature_para = text_frame.add_paragraph()
    feature_para.level = 0
    feature_para.bullet = False

    feature_bullet_run = feature_para.add_run()
    feature_bullet_run.text = "▪ "
    feature_bullet_run.font.color.rgb = state_colors.get(feature_state, RGBColor(0, 0, 0))
    feature_bullet_run.font.name = 'Segoe UI Light'

    run = feature_para.add_run()
    run.text = f"[Feature {feature_id}] - {feature_title}"
    run.font.bold = True
    run.font.name = 'Segoe UI Light'

def create_user_story_row(story, text_frame):

    story_para = text_frame.add_paragraph()
    story_para.level = 1
    story_para.bullet = False  # Disable bullet
    story_para.left_indent = Inches(0.5)  # Indentation for sub-level
    # Set text color based on story state
    story_state = story.get("state", "")
    
    story_bullet_run = story_para.add_run()
    story_bullet_run.text = "▪ "
    story_bullet_run.font.color.rgb = state_colors.get(story_state, RGBColor(0, 0, 0))
    story_bullet_run.font.name = 'Segoe UI Light'
    
    story_id_run = story_para.add_run()
    story_id_run.text = f"US {story['id']}: "
    story_id_run.font.bold = True
    story_id_run.font.name = 'Segoe UI Light'

    story_title_run = story_para.add_run()
    story_title_run.text = f"{story['title']}"
    story_title_run.font.bold = False
    story_title_run.font.name = 'Segoe UI Light'

def update_presentation_with_user_stories(user_stories, template_path, slide_index):
    
    prs = Presentation(template_path)
    slide = prs.slides[slide_index]

    # Look for BODY placeholder
    body_placeholder = None
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY:
            body_placeholder = placeholder
            break

    if not body_placeholder:
       raise ValueError("No BODY placeholder found on the specified slide.")

    text_frame = body_placeholder.text_frame
    text_frame.clear()

    # Group stories by feature
    stories_by_feature = defaultdict(list)
    for story in user_stories:
        feature = story['feature']
        feature_id = feature['id']
        stories_by_feature[feature_id, feature['title'], feature['state']].append(story)

    for (feature_id, feature_title, feature_state), stories in stories_by_feature.items():
        create_feature_row(feature_id, feature_title, feature_state, text_frame)

        # Add user stories under this feature as second-level bullets (level 1)
        for story in stories:
            create_user_story_row(story, text_frame)

    filename = f"MeetingMinutes-{datetime.today().strftime('%d%m%Y')}-SprintReview.pptx"
    prs.save(filename)
    print(f"Saved: {filename}")


# Get work items in the iteration (Sprint)
url = f"https://tfs-product.cmf.criticalmanufacturing.com/Products/{project}/{team}/_apis/work/teamsettings/iterations?$timeframe=current&api-version=7.0"
print('URL ', url)
iteration_res = requests.get(url, headers=headers, verify=False).json()

# (Optional) Log the iterations to make sure we’re targeting the right one
iteration_id = iteration_res["value"][0]["id"]

# Get work items in that iteration
# Note: May need to hardcode the iteration ID if this doesn’t return the sprint you want
backlog_url = f"https://tfs-product.cmf.criticalmanufacturing.com/Products/{project}/{team}/_apis/work/teamsettings/iterations/{iteration_id}/workitems?api-version=7.0"
backlog_res = requests.get(backlog_url, headers=headers, verify=False).json()

# List to hold only User Story IDs
user_story_ids = []

# Iterate through the work item relations and extract target IDs
for item in backlog_res['workItemRelations']:
    target_id = item['target']['id']
    
    # Fetch work item details (assuming a function to get work item details by ID)
    work_item_type = get_work_item_type(target_id)  # Assuming this function is defined to fetch work item type
    
    # Check if the work item is a User Story
    if work_item_type == 'User Story':
        user_story_ids.append(target_id)

# Output the User Story IDs

user_stories = get_work_item_details_w_features(user_story_ids)

update_presentation_with_user_stories(user_stories, template_path, 10)


